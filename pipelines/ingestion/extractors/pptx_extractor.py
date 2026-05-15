from pptx import Presentation

from pipelines.ingestion.models import ExtractedPage


async def extract_pptx(file_path: str):

    prs = Presentation(file_path)

    pages = []

    for i, slide in enumerate(prs.slides):

        text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text.append(shape.text)

        combined = "\n".join(text)

        pages.append(
            ExtractedPage(
                page_number=i + 1,
                content=combined,
                content_type="text",
                metadata={
                    "source": "pptx"
                }
            )
        )

    return pages